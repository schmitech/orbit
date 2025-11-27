from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

prs = Presentation()

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Sample PowerPoint for python-pptx Testing"
slide.placeholders[1].text = "Generated locally • Compatible with python-pptx 1.0.2"

# Slide 2: Bullets
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Agenda"
tf = slide.placeholders[1].text_frame
tf.clear()
agenda = [
    "Purpose of this deck",
    "Basic shapes and text",
    "Tables and charts",
    "Images and notes",
    "Summary slide",
]
for i, item in enumerate(agenda):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = item
    p.level = 0

# Slide 3: Shapes + Textbox
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Shapes and Text Boxes"

rect = slide.shapes.add_shape(
    1, Inches(1), Inches(1.8), Inches(3), Inches(1.2)  # 1 = rectangle
)
rect.text = "Rectangle Shape"
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor(220, 235, 250)
rect.line.color.rgb = RGBColor(60, 90, 140)

textbox = slide.shapes.add_textbox(Inches(4.5), Inches(1.8), Inches(4), Inches(1.5))
tb_tf = textbox.text_frame
tb_tf.text = "Text box with\nmultiple lines\nand styling."
tb_tf.paragraphs[0].font.size = Pt(18)
tb_tf.paragraphs[0].font.bold = True

# Slide 4: Table
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Sample Table"

table_shape = slide.shapes.add_table(
    4, 3, Inches(1), Inches(1.8), Inches(8), Inches(1.8)
)
table = table_shape.table

headers = ["Item", "Qty", "Cost (CAD)"]
for c, h in enumerate(headers):
    cell = table.cell(0, c)
    cell.text = h
    cell.text_frame.paragraphs[0].font.bold = True

data = [
    ("Compute", "2", "120.00"),
    ("Storage", "5 TB", "75.50"),
    ("Monitoring", "1", "19.99"),
]
for r, row in enumerate(data, start=1):
    for c, val in enumerate(row):
        table.cell(r, c).text = val

# Slide 5: Chart
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Simple Chart"

chart_data = ChartData()
chart_data.categories = ["Jan", "Feb", "Mar", "Apr"]
chart_data.add_series("Spend", (1200, 1500, 900, 1700))

slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1), Inches(1.8), Inches(8), Inches(3),
    chart_data
)

# Slide 6: Image placeholder + Notes
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Image and Speaker Notes"

ph = slide.shapes.add_shape(1, Inches(1), Inches(1.8), Inches(4), Inches(3))
ph.fill.solid()
ph.fill.fore_color.rgb = RGBColor(240, 240, 240)
ph.line.color.rgb = RGBColor(180, 180, 180)
ph.text = "Image Placeholder"

slide.notes_slide.notes_text_frame.text = (
    "These are sample speaker notes for testing note extraction."
)

# Slide 7: Summary
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Summary"
tf = slide.placeholders[1].text_frame
tf.clear()
summary = [
    "Deck includes titles, bullets, shapes, a table, a chart, and notes.",
    "Useful for testing parsing, ordering, and content extraction.",
    "Created to be small but feature-complete."
]
for i, item in enumerate(summary):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = item
    p.level = 0

prs.save("sample_python_pptx_test_deck.pptx")
print("Wrote sample_python_pptx_test_deck.pptx")