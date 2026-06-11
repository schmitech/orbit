"""
Document Generation Service

Renders a structured JSON document specification into binary document bytes using
native Python libraries. No external API calls — all rendering is local.

Supported formats: pdf (reportlab), docx (python-docx), xlsx (openpyxl), pptx (python-pptx).
These libraries are available in the 'files' dependency profile.

Spec schema expected from the LLM:
  {
    "title": str,
    "sections": [
      {
        "heading": str,
        "body": str,                    # optional paragraph text
        "table": [[str, ...]],          # optional; first row = header
        "bullet_points": [str, ...]     # optional
      }
    ],
    "metadata": {"author": str, "date": str}
  }
"""

import io
import logging
from typing import Dict, Any, Optional

logger = logging.getLogger(__name__)

MIME_TYPES: Dict[str, str] = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


class DocumentRenderer:
    """Render a document spec dict to bytes in the requested format.

    Pass the ``document_renderer`` section of ``document.yaml`` as ``config`` to
    override colours, fonts, margins, and other layout defaults.  Omitting ``config``
    (or any key within it) keeps the built-in defaults.
    """

    def __init__(self, config: Optional[Dict[str, Any]] = None) -> None:
        self._cfg: Dict[str, Any] = config or {}

    # ------------------------------------------------------------------
    # Config helper
    # ------------------------------------------------------------------

    def _get(self, fmt: str, *keys, default=None):
        """Safe nested lookup: self._cfg[fmt][keys[0]][keys[1]]...

        Returns *default* if any key is missing or the node is not a dict.
        An empty dict (key present but no sub-keys) is also treated as missing.
        """
        node = self._cfg.get(fmt, {})
        for k in keys:
            if not isinstance(node, dict):
                return default
            node = node.get(k)
            if node is None:
                return default
        if node == {}:
            return default
        return node

    def _meta(self, key: str, default: str = "") -> str:
        return self._cfg.get("metadata", {}).get(key, default)

    def render(self, spec: Dict[str, Any], fmt: str) -> bytes:
        fmt = fmt.lower()
        if fmt == "pdf":
            return self._render_pdf(spec)
        if fmt == "docx":
            return self._render_docx(spec)
        if fmt == "xlsx":
            return self._render_xlsx(spec)
        if fmt == "pptx":
            return self._render_pptx(spec)
        raise ValueError(f"Unsupported document format: {fmt!r}. Supported: pdf, docx, xlsx, pptx")

    # ------------------------------------------------------------------
    # Shared helpers
    # ------------------------------------------------------------------

    def _build_meta_line(self, spec: Dict[str, Any]) -> str:
        """Return the "Author: X  |  Date: Y" meta string for the document header."""
        meta = spec.get("metadata", {})
        parts = []
        # Configured author always wins — the LLM value is only a fallback when
        # no author is set in document.yaml.
        configured_author = self._meta("author", "")
        author = configured_author or meta.get("author", "ORBIT")
        org = self._meta("organization", "")
        display_author = f"{author} — {org}" if org else author
        parts.append(f"Author: {display_author}")
        if meta.get("date"):
            parts.append(f"Date: {meta['date']}")
        return "  |  ".join(parts)

    # ------------------------------------------------------------------
    # PDF — reportlab
    # ------------------------------------------------------------------

    def _render_pdf(self, spec: Dict[str, Any]) -> bytes:
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib import colors

        top_margin = self._get('pdf', 'margins', 'top', default=2) * cm
        bottom_margin = self._get('pdf', 'margins', 'bottom', default=2) * cm
        left_margin = self._get('pdf', 'margins', 'left', default=2.5) * cm
        right_margin = self._get('pdf', 'margins', 'right', default=2.5) * cm

        pagesize = (
            landscape(A4)
            if self._pdf_should_use_landscape(spec, A4, left_margin)
            else A4
        )
        buf = io.BytesIO()
        doc = SimpleDocTemplate(
            buf, pagesize=pagesize,
            topMargin=top_margin, bottomMargin=bottom_margin,
            leftMargin=left_margin, rightMargin=right_margin,
        )
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            "DocTitle", parent=styles["Title"],
            fontSize=self._get('pdf', 'typography', 'title_size', default=18),
            spaceAfter=self._get('pdf', 'typography', 'title_space_after', default=12),
        )
        h1_style = ParagraphStyle(
            "DocH1", parent=styles["Heading1"],
            fontSize=self._get('pdf', 'typography', 'heading1_size', default=14),
            spaceAfter=self._get('pdf', 'typography', 'heading1_space_after', default=6),
        )
        bullet_style = ParagraphStyle(
            "DocBullet", parent=styles["BodyText"],
            leftIndent=self._get('pdf', 'typography', 'bullet_indent', default=20),
            bulletIndent=self._get('pdf', 'typography', 'bullet_marker_indent', default=10),
        )

        story = []
        story.append(Paragraph(spec.get("title", "Document"), title_style))

        meta_line = self._build_meta_line(spec)
        if meta_line:
            story.append(Paragraph(meta_line, styles["Italic"]))
        story.append(Spacer(1, 0.5 * cm))

        for section in spec.get("sections", []):
            if section.get("heading"):
                story.append(Paragraph(section["heading"], h1_style))
            if section.get("body"):
                story.append(Paragraph(section["body"], styles["BodyText"]))
                story.append(Spacer(1, 0.3 * cm))
            for point in section.get("bullet_points") or []:
                story.append(Paragraph(f"• {point}", bullet_style))
            if section.get("bullet_points"):
                story.append(Spacer(1, 0.3 * cm))
            table_data = section.get("table")
            if table_data and len(table_data) > 0:
                tbl = self._build_pdf_table(
                    table_data=table_data,
                    styles=styles,
                    available_width=doc.width,
                )
                story.append(tbl)
                story.append(Spacer(1, 0.3 * cm))

        doc.build(story)
        return buf.getvalue()

    def _pdf_should_use_landscape(self, spec: Dict[str, Any], portrait_pagesize, horizontal_margin: float) -> bool:
        portrait_width = portrait_pagesize[0] - (2 * horizontal_margin)
        threshold = self._get('pdf', 'table', 'landscape_threshold_cols', default=7)
        for section in spec.get("sections", []):
            table_data = section.get("table")
            if not table_data:
                continue
            normalized_rows = self._normalize_table_rows(table_data)
            if normalized_rows and len(normalized_rows[0]) >= threshold:
                return True
            if self._estimate_table_width(table_data) > portrait_width:
                return True
        return False

    def _build_pdf_table(self, table_data, styles, available_width: float):
        from reportlab.lib import colors
        from reportlab.lib.styles import ParagraphStyle
        from reportlab.pdfbase.pdfmetrics import stringWidth
        from reportlab.platypus import Paragraph, Table, TableStyle

        font_body = self._get('pdf', 'fonts', 'body', default='Helvetica')
        font_header = self._get('pdf', 'fonts', 'header', default='Helvetica-Bold')
        color_header_bg = self._get('pdf', 'colors', 'header_bg', default='#4472C4')
        color_header_text = self._get('pdf', 'colors', 'header_text', default='white')
        color_alt_row = self._get('pdf', 'colors', 'alt_row_bg', default='#EBF0FA')
        color_grid = self._get('pdf', 'colors', 'grid', default='grey')
        cell_padding = self._get('pdf', 'table', 'cell_padding', default=8)

        normalized_rows = self._normalize_table_rows(table_data)
        column_count = len(normalized_rows[0]) if normalized_rows else 0
        font_size = self._resolve_table_font_size(column_count)

        body_style = ParagraphStyle(
            "DocTableBody",
            parent=styles["BodyText"],
            fontName=font_body,
            fontSize=font_size,
            leading=font_size + 2,
            wordWrap="CJK",
        )
        header_text_color = (
            colors.HexColor(color_header_text)
            if color_header_text.startswith('#')
            else getattr(colors, color_header_text, colors.white)
        )
        header_style = ParagraphStyle(
            "DocTableHeader",
            parent=body_style,
            fontName=font_header,
            textColor=header_text_color,
        )

        col_widths = self._compute_pdf_col_widths(
            normalized_rows=normalized_rows,
            available_width=available_width,
            font_size=font_size,
            string_width=stringWidth,
            font_body=font_body,
            font_header=font_header,
            cell_padding=cell_padding,
        )
        wrapped_rows = [
            [
                Paragraph(str(cell), header_style if row_idx == 0 else body_style)
                for cell in row
            ]
            for row_idx, row in enumerate(normalized_rows)
        ]

        tbl = Table(
            wrapped_rows,
            colWidths=col_widths,
            repeatRows=1,
            splitByRow=1,
            hAlign="LEFT",
        )
        grid_color = (
            colors.HexColor(color_grid)
            if color_grid.startswith('#')
            else getattr(colors, color_grid, colors.grey)
        )
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(color_header_bg)),
            ("TEXTCOLOR", (0, 0), (-1, 0), header_text_color),
            ("FONTNAME", (0, 0), (-1, 0), font_header),
            ("FONTSIZE", (0, 0), (-1, -1), font_size),
            ("GRID", (0, 0), (-1, -1), 0.5, grid_color),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor(color_alt_row)]),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), cell_padding / 2),
            ("RIGHTPADDING", (0, 0), (-1, -1), cell_padding / 2),
            ("TOPPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ]))
        return tbl

    def _estimate_table_width(self, table_data) -> float:
        from reportlab.pdfbase.pdfmetrics import stringWidth

        font_body = self._get('pdf', 'fonts', 'body', default='Helvetica')
        font_header = self._get('pdf', 'fonts', 'header', default='Helvetica-Bold')
        cell_padding = self._get('pdf', 'table', 'cell_padding', default=8)
        font_size = self._get('pdf', 'table', 'font_size', default=9)

        normalized_rows = self._normalize_table_rows(table_data)
        widths = self._compute_pdf_col_widths(
            normalized_rows=normalized_rows,
            available_width=None,
            font_size=font_size,
            string_width=stringWidth,
            font_body=font_body,
            font_header=font_header,
            cell_padding=cell_padding,
        )
        return sum(widths)

    @staticmethod
    def _normalize_table_rows(table_data):
        safe_rows = [[str(cell) for cell in row] for row in table_data if row]
        if not safe_rows:
            return []
        column_count = max(len(row) for row in safe_rows)
        return [row + [""] * (column_count - len(row)) for row in safe_rows]

    def _resolve_table_font_size(self, column_count: int) -> int:
        if column_count >= 8:
            return self._get('pdf', 'table', 'min_font_size', default=7)
        if column_count >= 6:
            return self._get('pdf', 'table', 'font_size_many_cols', default=8)
        return self._get('pdf', 'table', 'font_size', default=9)

    def _compute_pdf_col_widths(self, normalized_rows, available_width, font_size: int, string_width,
                                font_body: str = 'Helvetica', font_header: str = 'Helvetica-Bold',
                                cell_padding: int = 8):
        if not normalized_rows:
            return []

        column_count = len(normalized_rows[0])
        widths = []
        min_width = self._get('pdf', 'table', 'min_col_width', default=42)
        max_col_ratio = self._get('pdf', 'table', 'max_col_ratio', default=0.28)
        max_width = available_width * max_col_ratio if available_width else None

        for col_idx in range(column_count):
            widest = 0.0
            for row_idx, row in enumerate(normalized_rows):
                font_name = font_header if row_idx == 0 else font_body
                cell_text = row[col_idx]
                widest = max(widest, string_width(cell_text, font_name, font_size))
            padded = widest + cell_padding
            if max_width is not None:
                padded = min(padded, max_width)
            widths.append(max(min_width, padded))

        if available_width and sum(widths) > available_width:
            scale = available_width / sum(widths)
            widths = [max(min_width, width * scale) for width in widths]

            if sum(widths) > available_width:
                even_width = available_width / column_count
                widths = [even_width] * column_count
            else:
                widths[-1] += available_width - sum(widths)

        return widths

    # ------------------------------------------------------------------
    # DOCX — python-docx
    # ------------------------------------------------------------------

    def _render_docx(self, spec: Dict[str, Any]) -> bytes:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()
        title_para = doc.add_heading(spec.get("title", "Document"), level=0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

        meta_line = self._build_meta_line(spec)
        if meta_line:
            p = doc.add_paragraph()
            run = p.add_run(meta_line)
            run.italic = True
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER

        doc.add_paragraph()

        table_style = self._get('docx', 'table_style', default='Table Grid')

        for section in spec.get("sections", []):
            if section.get("heading"):
                doc.add_heading(section["heading"], level=1)
            if section.get("body"):
                doc.add_paragraph(section["body"])
            for point in section.get("bullet_points") or []:
                doc.add_paragraph(point, style="List Bullet")
            table_data = section.get("table")
            if table_data and len(table_data) > 1:
                headers = table_data[0]
                rows = table_data[1:]
                tbl = doc.add_table(rows=1 + len(rows), cols=len(headers))
                tbl.style = table_style
                hdr_cells = tbl.rows[0].cells
                for i, h in enumerate(headers):
                    hdr_cells[i].text = str(h)
                    hdr_cells[i].paragraphs[0].runs[0].bold = True
                for r_idx, row in enumerate(rows):
                    for c_idx, val in enumerate(row):
                        tbl.rows[r_idx + 1].cells[c_idx].text = str(val)
                doc.add_paragraph()

        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()

    # ------------------------------------------------------------------
    # XLSX — openpyxl
    # ------------------------------------------------------------------

    def _render_xlsx(self, spec: Dict[str, Any]) -> bytes:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment

        header_color = self._get('xlsx', 'colors', 'header_bg', default='4472C4')
        alt_color = self._get('xlsx', 'colors', 'alt_row_bg', default='EBF0FA')
        summary_col_width = self._get('xlsx', 'summary_col_width', default=80)
        max_col_width = self._get('xlsx', 'max_col_width', default=50)

        wb = Workbook()
        wb.remove(wb.active)  # Remove the default blank sheet

        header_fill = PatternFill(start_color=header_color, end_color=header_color, fill_type="solid")
        alt_fill = PatternFill(start_color=alt_color, end_color=alt_color, fill_type="solid")

        # Summary sheet
        ws = wb.create_sheet("Summary")
        ws["A1"] = spec.get("title", "Document")
        ws["A1"].font = Font(size=16, bold=True)

        meta_line = self._build_meta_line(spec)
        if meta_line:
            ws["A2"] = meta_line
            ws["A2"].font = Font(italic=True)
        ws.column_dimensions["A"].width = summary_col_width

        row = 4
        for section in spec.get("sections", []):
            if section.get("heading"):
                ws.cell(row=row, column=1, value=section["heading"]).font = Font(bold=True, size=12)
                row += 1
            if section.get("body"):
                ws.cell(row=row, column=1, value=section["body"])
                row += 1
            for point in section.get("bullet_points") or []:
                ws.cell(row=row, column=1, value=f"• {point}")
                row += 1
            row += 1

        # Data sheets — one per section that has a table
        for idx, section in enumerate(spec.get("sections", [])):
            table_data = section.get("table")
            if not table_data or len(table_data) < 1:
                continue
            sheet_name = (section.get("heading") or f"Sheet{idx + 1}")[:31]
            ws_data = wb.create_sheet(sheet_name)
            for r_idx, tbl_row in enumerate(table_data):
                for c_idx, val in enumerate(tbl_row):
                    cell = ws_data.cell(row=r_idx + 1, column=c_idx + 1, value=str(val))
                    if r_idx == 0:
                        cell.font = Font(bold=True, color="FFFFFF")
                        cell.fill = header_fill
                        cell.alignment = Alignment(horizontal="center")
                    elif r_idx % 2 == 0:
                        cell.fill = alt_fill
            # Auto-fit column widths
            for col in ws_data.columns:
                max_len = max((len(str(c.value or "")) for c in col), default=10)
                ws_data.column_dimensions[col[0].column_letter].width = min(max_len + 4, max_col_width)

        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue()

    # ------------------------------------------------------------------
    # PPTX — python-pptx
    # ------------------------------------------------------------------

    def _render_pptx(self, spec: Dict[str, Any]) -> bytes:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        slide_width = self._get('pptx', 'slide_width_inches', default=13.33)
        slide_height = self._get('pptx', 'slide_height_inches', default=7.5)
        body_pt = self._get('pptx', 'typography', 'body_pt', default=14)
        bullet_pt = self._get('pptx', 'typography', 'bullet_pt', default=13)
        row_height = self._get('pptx', 'table', 'row_height_inches', default=0.4)

        prs = Presentation()
        prs.slide_width = Inches(slide_width)
        prs.slide_height = Inches(slide_height)

        # Title slide
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = spec.get("title", "Document")
        meta_line = self._build_meta_line(spec)
        if len(slide.placeholders) > 1 and meta_line:
            slide.placeholders[1].text = meta_line

        # Content slides — one per section
        content_layout = prs.slide_layouts[1]
        blank_layout = prs.slide_layouts[6]
        for section in spec.get("sections", []):
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = section.get("heading", "")
            tf = slide.placeholders[1].text_frame
            tf.clear()

            body = section.get("body", "")
            if body:
                p = tf.paragraphs[0]
                p.text = body
                p.font.size = Pt(body_pt)

            for point in section.get("bullet_points") or []:
                p = tf.add_paragraph()
                p.text = point
                p.level = 1
                p.font.size = Pt(bullet_pt)

            # Tables go on their own slide
            table_data = section.get("table")
            if table_data and len(table_data) > 1:
                tbl_slide = prs.slides.add_slide(blank_layout)
                rows = len(table_data)
                cols = len(table_data[0])
                tbl = tbl_slide.shapes.add_table(
                    rows, cols,
                    Inches(0.5), Inches(1.5),
                    Inches(slide_width - 1),
                    Inches(min(rows * row_height + 0.5, 5.5)),
                ).table
                for r, row in enumerate(table_data):
                    for c, val in enumerate(row):
                        cell = tbl.cell(r, c)
                        cell.text = str(val)
                        if r == 0:
                            cell.text_frame.paragraphs[0].font.bold = True
                            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()
