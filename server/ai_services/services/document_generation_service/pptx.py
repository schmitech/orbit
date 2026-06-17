import io
from typing import Dict, Any

from .base import BaseRenderer


class PptxRenderer(BaseRenderer):
    """Render a document spec to PPTX bytes using python-pptx."""

    @staticmethod
    def _hex_to_rgb(hex_str: str):
        from pptx.dml.color import RGBColor
        h = hex_str.lstrip('#')
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def render(self, spec: Dict[str, Any]) -> bytes:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        slide_width   = self._get('pptx', 'slide_width_inches',  default=13.33)
        slide_height  = self._get('pptx', 'slide_height_inches', default=7.5)
        color_hdr_bg  = self._get('pptx', 'colors', 'header_bg',   default='4472C4')
        color_hdr_txt = self._get('pptx', 'colors', 'header_text', default='FFFFFF')
        color_alt_row = self._get('pptx', 'colors', 'alt_row_bg',  default='EBF0FA')
        font_body     = self._get('pptx', 'fonts',  'body',         default='Calibri')
        font_header   = self._get('pptx', 'fonts',  'header',       default='Calibri')
        title_pt      = self._get('pptx', 'typography', 'title_pt',  default=28)
        body_pt       = self._get('pptx', 'typography', 'body_pt',   default=14)
        bullet_pt     = self._get('pptx', 'typography', 'bullet_pt', default=13)
        row_height    = self._get('pptx', 'table', 'row_height_inches', default=0.4)
        show_heading  = self._get('pptx', 'table', 'show_heading',      default=True)
        slide_numbers = self._get('pptx', 'slide_numbers',              default=False)

        rgb_hdr_bg  = self._hex_to_rgb(color_hdr_bg)
        rgb_hdr_txt = self._hex_to_rgb(color_hdr_txt)
        rgb_alt_row = self._hex_to_rgb(color_alt_row)

        prs = Presentation()
        prs.slide_width  = Inches(slide_width)
        prs.slide_height = Inches(slide_height)

        # Title slide
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        title_shape = slide.shapes.title
        title_shape.text = spec.get("title", "Document")
        if title_shape.text_frame.paragraphs:
            run = title_shape.text_frame.paragraphs[0].runs
            if run:
                run[0].font.name = font_header
        meta_line = self._build_meta_line(spec)
        if len(slide.placeholders) > 1 and meta_line:
            ph = slide.placeholders[1]
            ph.text = meta_line
            if ph.text_frame.paragraphs:
                r = ph.text_frame.paragraphs[0].runs
                if r:
                    r[0].font.name = font_body
                    r[0].font.size = Pt(body_pt)

        content_layout = prs.slide_layouts[1]
        blank_layout   = prs.slide_layouts[6]
        slide_index = 1

        for section in spec.get("sections", []):
            heading = section.get("heading", "")
            body    = section.get("body", "")
            bullets = section.get("bullet_points") or []

            if body or bullets:
                slide = prs.slides.add_slide(content_layout)
                slide_index += 1
                title_tf = slide.shapes.title.text_frame
                title_tf.text = heading
                if title_tf.paragraphs and title_tf.paragraphs[0].runs:
                    title_tf.paragraphs[0].runs[0].font.name = font_header
                    title_tf.paragraphs[0].runs[0].font.size = Pt(title_pt)

                tf = slide.placeholders[1].text_frame
                tf.clear()
                if body:
                    p = tf.paragraphs[0]
                    p.text = body
                    if p.runs:
                        p.runs[0].font.name = font_body
                    p.font.size = Pt(body_pt)
                for point in bullets:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 1
                    if p.runs:
                        p.runs[0].font.name = font_body
                    p.font.size = Pt(bullet_pt)

                if slide_numbers:
                    self._add_slide_number(slide, slide_index, slide_width, slide_height, font_body)

            table_data = section.get("table")
            if table_data and len(table_data) > 1:
                tbl_slide = prs.slides.add_slide(blank_layout)
                slide_index += 1
                rows = len(table_data)
                cols = len(table_data[0])

                top_offset = Inches(0.6)
                if show_heading and heading:
                    lbl = tbl_slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.2),
                        Inches(slide_width - 1), Inches(0.5),
                    )
                    lbl_tf = lbl.text_frame
                    lbl_tf.text = heading
                    lbl_p = lbl_tf.paragraphs[0]
                    lbl_p.font.bold = True
                    lbl_p.font.size = Pt(title_pt)
                    lbl_p.font.name = font_header
                    top_offset = Inches(0.9)

                tbl_height = Inches(min(rows * row_height + 0.3, slide_height - top_offset / 914400))
                tbl = tbl_slide.shapes.add_table(
                    rows, cols,
                    Inches(0.5), top_offset,
                    Inches(slide_width - 1), tbl_height,
                ).table

                for r, row in enumerate(table_data):
                    for c, val in enumerate(row):
                        cell = tbl.cell(r, c)
                        cell.text = str(val)
                        para = cell.text_frame.paragraphs[0]
                        if r == 0:
                            para.font.bold  = True
                            para.font.name  = font_header
                            para.font.color.rgb = rgb_hdr_txt
                            para.alignment = PP_ALIGN.CENTER
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = rgb_hdr_bg
                        else:
                            para.font.name = font_body
                            para.font.size = Pt(body_pt)
                            if r % 2 == 0:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = rgb_alt_row

                if slide_numbers:
                    self._add_slide_number(tbl_slide, slide_index, slide_width, slide_height, font_body)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    @staticmethod
    def _add_slide_number(slide, number: int, slide_width: float, slide_height: float, font_name: str) -> None:
        from pptx.util import Inches, Pt
        tb = slide.shapes.add_textbox(
            Inches(slide_width - 0.6), Inches(slide_height - 0.35),
            Inches(0.5), Inches(0.3),
        )
        p = tb.text_frame.paragraphs[0]
        p.text = str(number)
        p.font.size = Pt(10)
        p.font.name = font_name
