"""
PPTX Processor

Handles PowerPoint PPTX files using python-pptx.
"""

import logging
from typing import Dict, Any
from .base_processor import FileProcessor

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from io import BytesIO
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available. PPTX processing disabled.")


class PPTXProcessor(FileProcessor):
    """
    Processor for PowerPoint PPTX files.

    Supports: application/vnd.openxmlformats-officedocument.presentationml.presentation
    Requires: python-pptx
    """

    def supports_mime_type(self, mime_type: str) -> bool:
        """Check if this processor supports the MIME type."""
        pptx_types = [
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'application/vnd.ms-powerpoint',
        ]
        return PPTX_AVAILABLE and mime_type.lower() in pptx_types

    async def extract_text(self, file_data: bytes, filename: str = None) -> str:
        """Extract text from PPTX."""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not available")

        logger.debug(f"PPTXProcessor.extract_text() called for file: {filename or 'unknown'} (using python-pptx)")

        text_parts = []

        try:
            prs = Presentation(BytesIO(file_data))

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []

                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)

                    # Extract text from tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_texts = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_texts.append(cell.text.strip())
                            if row_texts:
                                slide_texts.append(" | ".join(row_texts))

                if slide_texts:
                    text_parts.append(f"--- Slide {slide_num} ---")
                    text_parts.extend(slide_texts)
                    text_parts.append("")  # Empty line between slides

            return "\n".join(text_parts)

        except Exception as e:
            logger.error(f"Error processing PPTX: {e}")
            raise

    async def extract_metadata(self, file_data: bytes, filename: str = None) -> Dict[str, Any]:
        """Extract metadata from PPTX."""
        metadata = await super().extract_metadata(file_data, filename)

        if not PPTX_AVAILABLE:
            return metadata

        try:
            prs = Presentation(BytesIO(file_data))

            # Count slides
            slide_count = len(prs.slides)

            # Get core properties if available
            core_props = prs.core_properties

            metadata.update({
                'slide_count': slide_count,
                'mime_type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            })

            # Add core properties if available
            if core_props.title:
                metadata['title'] = core_props.title
            if core_props.author:
                metadata['author'] = core_props.author
            if core_props.subject:
                metadata['subject'] = core_props.subject
            if core_props.created:
                metadata['created'] = core_props.created.isoformat() if core_props.created else None
            if core_props.modified:
                metadata['modified'] = core_props.modified.isoformat() if core_props.modified else None

        except Exception as e:
            logger.warning(f"Error extracting PPTX metadata: {e}")

        return metadata
